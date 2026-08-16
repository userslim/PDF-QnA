"""导出模块 - 生成 PPT 和 PDF"""
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT
from datetime import datetime


class ExportManager:
    """导出 PPT 和 PDF"""
    
    def export_to_ppt(
        self, 
        session_name: str,
        qa_history: List[Dict],
        output_path: str
    ) -> str:
        """导出问答记录为 PPT"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = session_name
        subtitle.text = f"文档问答总结\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # 为每个问答创建幻灯片
        for i, qa in enumerate(qa_history, 1):
            # 问题页
            question_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(question_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"问题 {i}: {qa['question'][:50]}"
            
            tf = content.text_frame
            tf.text = qa['answer']
            
            # 设置文本格式
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.LEFT
            
            # 添加来源页
            if qa.get('sources'):
                source_slide = prs.slides.add_slide(prs.slide_layouts[1])
                source_slide.shapes.title.text = f"参考资料 - 问题 {i}"
                
                tf = source_slide.placeholders[1].text_frame
                tf.text = ""
                
                for j, source in enumerate(qa['sources'][:5], 1):
                    p = tf.add_paragraph()
                    p.text = f"{j}. {source['source']} - 页码 {source['page']}"
                    p.level = 0
                    p.font.size = Pt(12)
        
        prs.save(output_path)
        return output_path
    
    def export_to_pdf(
        self,
        session_name: str,
        qa_history: List[Dict],
        output_path: str
    ) -> str:
        """导出问答记录为 PDF"""
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # 自定义样式
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor='#1f4e79',
            spaceAfter=30,
            alignment=TA_LEFT
        )
        
        question_style = ParagraphStyle(
            'Question',
            parent=styles['Heading2'],
            fontSize=16,
            textColor='#2e75b6',
            spaceAfter=12,
            spaceBefore=20
        )
        
        answer_style = ParagraphStyle(
            'Answer',
            parent=styles['BodyText'],
            fontSize=11,
            leftIndent=20,
            spaceAfter=15
        )
        
        source_style = ParagraphStyle(
            'Source',
            parent=styles['BodyText'],
            fontSize=9,
            textColor='#666666',
            leftIndent=40,
            spaceAfter=5
        )
        
        # 标题
        story.append(Paragraph(session_name, title_style))
        story.append(Paragraph(
            f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            styles['Normal']
        ))
        story.append(Spacer(1, 0.3 * inch))
        
        # 问答内容
        for i, qa in enumerate(qa_history, 1):
            story.append(Paragraph(f"问题 {i}: {qa['question']}", question_style))
            story.append(Paragraph(qa['answer'], answer_style))
            
            # 来源
            if qa.get('sources'):
                story.append(Paragraph("参考资料:", styles['Heading4']))
                for j, source in enumerate(qa['sources'][:5], 1):
                    src_text = f"[{j}] {source['source']} - 页码 {source['page']}"
                    story.append(Paragraph(src_text, source_style))
            
            story.append(PageBreak())
        
        doc.build(story)
        return output_path


if __name__ == "__main__":
    em = ExportManager()
    print("导出管理器初始化成功")